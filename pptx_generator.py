import json
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

from langchain_openai import ChatOpenAI

from langchain_core.tools import tool
from langchain_community.tools import DuckDuckGoSearchRun
from langchain_community.utilities.dalle_image_generator import DallEAPIWrapper
from langchain.agents import create_tool_calling_agent, AgentExecutor
from langchain_core.prompts import ChatPromptTemplate

from dotenv import load_dotenv
load_dotenv()

# Tool 1: Web Search
search_tool = DuckDuckGoSearchRun(
    name="web_search",
    description="Use this to search the web for recent facts and statistics about the topic."
)

# Tool 2: Image Generation (Using DALL-E)
@tool
def generate_slide_image(prompt: str) -> str:
    """Generates an image based on a prompt and returns the image URL. Use this when a slide needs a visual."""
    try:
        dalle = DallEAPIWrapper(size="1024x1024", model="dall-e-3")
        image_url = dalle.run(prompt)
        return image_url
    except Exception as e:
        return f"Image generation failed: {e}"

# Tool 3: PPTX Compiler
@tool
def create_pptx(slide_data_json: str) -> str:
    """
    Call this tool ONLY when you have researched the topic, generated image URLs, and finalized the content.
    Input MUST be a valid JSON string with this format:
    [{"title": "Slide 1", "bullets": ["point 1", "point 2"], "image_url": "url_or_null"}]
    """
    try:
        slide_data = json.loads(slide_data_json)
        prs = Presentation()
        
        for slide_info in slide_data:
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)
            
            if "title" in slide_info:
                slide.shapes.title.text = slide_info["title"]
                
            if "bullets" in slide_info:
                tf = slide.shapes.placeholders[1].text_frame
                for point in slide_info["bullets"]:
                    p = tf.add_paragraph()
                    p.text = point
                    
            image_url = slide_info.get("image_url")
            if image_url and image_url != "null":
                try:
                    response = requests.get(image_url)
                    image_bytes = BytesIO(response.content)
                    slide.shapes.add_picture(image_bytes, Inches(5.5), Inches(2), width=Inches(4))
                except Exception as e:
                    print(f"Skipping image due to error: {e}")
                    
        filename = "AI_Generated_Presentation.pptx"
        prs.save(filename)
        return f"SUCCESS! Presentation saved successfully as {filename}"
        
    except json.JSONDecodeError:
        return "ERROR: The input was not valid JSON. Please try again."

tools = [search_tool, generate_slide_image, create_pptx]

llm = ChatOpenAI(model="gpt-4o", temperature=0.3)

# Create the Agent Prompt
prompt = ChatPromptTemplate.from_messages([
    ("system", """You are an expert Presentation Creation Agent. 
    Your goal is to build a high-quality PowerPoint presentation for the user.
    
    Follow these steps strictly:
    1. Use the 'web_search' tool to gather recent, accurate facts about the user's topic.
    2. Draft an outline with 3-4 slides. Each slide needs a title and 3 solid bullet points.
    3. Use the 'generate_slide_image' tool to create 1 or 2 relevant images for the presentation. Save the returned URLs.
    4. Compile everything into the required JSON format.
    5. Call the 'create_pptx' tool passing the JSON string to build the final file.
    
    Do not stop until the 'create_pptx' tool returns a SUCCESS message.
    """),
    ("human", "{input}"),
    ("placeholder", "{agent_scratchpad}"),
])

agent = create_tool_calling_agent(llm, tools, prompt)
agent_executor = AgentExecutor(agent=agent, tools=tools, verbose=True)

#Running the agent

if __name__ == "__main__":
    user_topic = "The Impact of Quantum Computing on Cybersecurity"
    print(f"Starting Agent Run for topic: {user_topic}\n")
    
    response = agent_executor.invoke({"input": f"Create a presentation about: {user_topic}"})
    
    print(response["output"])